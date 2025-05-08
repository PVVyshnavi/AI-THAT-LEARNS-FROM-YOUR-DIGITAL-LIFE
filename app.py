import streamlit as st
import spacy
import nltk
nltk.download('punkt')

from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF for PDFs
import json

# Summarization function
def generate_summary(text, num_sentences=3):
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LsaSummarizer()
    summary = summarizer(parser.document, num_sentences)
    return " ".join(str(sentence) for sentence in summary)

# Extract text from supported file types
def extract_text_from_file(file):
    filename = file.name.lower()

    if filename.endswith(".txt"):
        return file.read().decode("utf-8")

    elif filename.endswith(".docx"):
        doc = Document(file)
        return "\n".join(p.text for p in doc.paragraphs)

    elif filename.endswith(".pptx"):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif filename.endswith(".pdf"):
        pdf_text = ""
        with fitz.open(stream=file.read(), filetype="pdf") as doc:
            for page in doc:
                pdf_text += page.get_text()
        return pdf_text

    else:
        return "Unsupported file format. Please upload .txt, .docx, .pptx, or .pdf files."

# Streamlit app configuration
st.set_page_config(page_title="AI That Learns From Your Digital Life", layout="wide")

# Load spaCy model
nlp = spacy.load("en_core_web_sm")

# App title
st.write("Upload `.txt`, `.docx`, `.pptx`, or `.pdf` files to extract insights from your personal notes.")

uploaded_files = st.file_uploader(
    "Upload supported files (.txt, .docx, .pptx, .pdf)", 
    type=["txt", "docx", "pptx", "pdf"], 
    accept_multiple_files=True
)

# ✅ File upload with proper label and supported types

if 'history' not in st.session_state:
    st.session_state.history = []

if uploaded_files:
    for uploaded_file in uploaded_files:
        text = extract_text_from_file(uploaded_file)

        st.subheader(f"📄 Uploaded File: {uploaded_file.name}")
        st.write(text[:1000] + "..." if len(text) > 1000 else text)

        if text.startswith("Unsupported file format"):
            st.warning(text)
            continue

        # NLP processing
        doc = nlp(text)

        # Keywords
        st.subheader("🔑 Keywords / Important Tokens:")
        keywords = [token.text for token in doc if token.is_alpha and not token.is_stop]
        top_keywords = list(set(keywords[:20]))
        st.write(top_keywords)

        st.markdown("---")

        # Word Cloud
        st.subheader("☁️ Word Cloud:")
        word_freq = " ".join(keywords)
        wordcloud = WordCloud(width=800, height=400, background_color='white').generate(word_freq)
        fig, ax = plt.subplots(figsize=(10, 5))
        ax.imshow(wordcloud, interpolation="bilinear")
        ax.axis("off")
        st.pyplot(fig)

        st.markdown("---")

        # Named Entities
        st.subheader("🏷 Named Entities:")
        for ent in doc.ents:
            st.write(f"{ent.text} — {ent.label_}")

        st.markdown("---")

        # Summary
        st.subheader("🧾 Summary:")
        summary = generate_summary(text, 3)
        st.write(summary)

        st.markdown("---")

        # Save to session history
        result = {
            "filename": uploaded_file.name,
            "keywords": top_keywords,
            "summary": summary,
            "entities": [(ent.text, ent.label_) for ent in doc.ents]
        }
        st.session_state.history.append(result)

        # Download options
        st.subheader("💾 Download Results:")
        json_data = json.dumps(result, indent=2)
        st.download_button("Download as JSON", data=json_data, file_name=f"{uploaded_file.name}_analysis.json", mime="application/json")

        txt_data = f"Summary:\n{summary}\n\nKeywords:\n{', '.join(top_keywords)}\n\nEntities:\n" + "\n".join(f"{ent[0]} — {ent[1]}" for ent in result['entities'])
        st.download_button("Download as TXT", data=txt_data, file_name=f"{uploaded_file.name}_analysis.txt", mime="text/plain")

    st.markdown("---")
    st.subheader("🧠 Session Memory (History)")
    for i, item in enumerate(st.session_state.history[::-1], 1):
        st.markdown(f"**{i}. {item['filename']}**")
        st.write(f"**Summary**: {item['summary']}")
